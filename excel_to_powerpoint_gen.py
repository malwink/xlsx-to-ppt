"""
Sample Excel to PowerPoint generator

Reads a workbook with:
- Sheet 1: settings
- Sheet 2: data table to dump into the final slide

Creates:
- Slide 1: Title-only slide styled from settings
- Slides 2..N-1: simple title/content slides
- Slide N: a "Final Results" slide containing an Excel table + a footnote
"""

from __future__ import annotations

import re
from dataclasses import dataclass
from typing import Optional, Tuple

import openpyxl
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


@dataclass(frozen=True)
class PptSettings:
    slide_count: int
    title_size_pt: int
    title_color: RGBColor
    text_size_pt: int


class ExcelToPPT:
    """
    Convert an Excel file into a PPTX deck using the first sheet as settings and
    the second sheet as table data.

    Excel expectations (settings sheet, row 2):
      A2: slide_count (int)
      B2: title_size in points (int)
      C2: title color in hex (#RRGGBB or RRGGBB)
      D2: body text size in points (int)
    """

    def __init__(self, excel_file: str, pptx_file: str):
        self.excel_file = excel_file
        self.pptx_file = pptx_file

        # data_only=True reads results of formulas
        self.workbook = openpyxl.load_workbook(self.excel_file, data_only=True)

        if len(self.workbook.worksheets) < 2:
            raise ValueError("Workbook must contain at least 2 sheets: settings + data.")

        self.settings_sheet = self.workbook.worksheets[0]
        self.data_sheet = self.workbook.worksheets[1]
        self.settings = self._read_settings()

    @staticmethod
    def _hex_to_rgb_tuple(hex_str: str) -> Tuple[int, int, int]:
        """
        Convert 'RRGGBB' to (R, G, B). Raises ValueError if it's malformed.
        """
        return tuple(int(hex_str[i : i + 2], 16) for i in (0, 2, 4))

    @staticmethod
    def _parse_hex_color(value: Optional[str], default: str = "1F4E79") -> RGBColor:
        """
        Accepts '#RRGGBB' or 'RRGGBB'. Falls back to default if missing/invalid.
        """
        if not value:
            rgb = ExcelToPPT._hex_to_rgb_tuple(default)
            return RGBColor(*rgb)

        # Strip whitespace and leading '#'
        s = str(value).strip()
        if s.startswith("#"):
            s = s[1:]

        # Validate hex
        if not re.fullmatch(r"[0-9a-fA-F]{6}", s):
            rgb = ExcelToPPT._hex_to_rgb_tuple(default)
            return RGBColor(*rgb)

        rgb = ExcelToPPT._hex_to_rgb_tuple(s)
        return RGBColor(*rgb)

    def _read_settings(self) -> PptSettings:
        """
        Reads and validates settings from the settings sheet.
        Provides reasonable defaults if some values are missing.
        """
        slide_count = int(self.settings_sheet["A2"].value)
        title_size = int(self.settings_sheet["B2"].value)
        title_color = self.settings_sheet["C2"].value
        text_size = int(self.settings_sheet["D2"].value)

        if slide_count < 2:
            raise ValueError("slide_count (A2) must be at least 2 (title slide + final slide).")

        return PptSettings(
            slide_count=slide_count,
            title_size_pt=title_size,
            title_color=self._parse_hex_color(str(title_color) if title_color else None),
            text_size_pt=text_size,
        )

    # ---------- PPT generation ----------

    def create_presentation(self) -> None:
        prs = Presentation()

        # 1) Title slide (layout 5 is typically "Title Only" in the default template)
        title_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = "Presentation"
        p = title_shape.text_frame.paragraphs[0]
        p.font.size = Pt(self.settings.title_size_pt)
        p.font.color.rgb = self.settings.title_color

        # 2) Middle slides
        # We create slides 2..(slide_count-1) as standard title/content slides.
        for i in range(2, self.settings.slide_count):
            slide_layout = prs.slide_layouts[1]
            s = prs.slides.add_slide(slide_layout)
            s.shapes.title.text = f"Slide {i - 1}"

            # Add some placeholder content so slides aren't empty
            try:
                body = s.placeholders[1].text_frame
                body.clear()
                body.text = "Content goes here…"
                body.paragraphs[0].font.size = Pt(self.settings.text_size_pt)
            except Exception:
                pass

        # 3) Final slide with table (always add this last)
        final_layout = prs.slide_layouts[1]
        final_slide = prs.slides.add_slide(final_layout)
        final_slide.shapes.title.text = "Final Results"

        self._add_table_from_sheet(final_slide)

        prs.save(self.pptx_file)

    def _add_table_from_sheet(self, slide) -> None:
        """
        Adds a PPTX table on the given slide + a footnote.
        """
        rows = self.data_sheet.max_row or 1
        cols = self.data_sheet.max_column or 1

        # Slide dimensions in inches
        slide_width_in = 10
        slide_height_in = 7.5

        margin_left = 0.7
        margin_right = 0.7
        top = 1.6  # below title
        bottom_margin = 1.0

        table_left = Inches(margin_left)
        table_top = Inches(top)
        table_width = Inches(slide_width_in - margin_left - margin_right)

        # Keep table from becoming ridiculously tall.
        max_table_height_in = slide_height_in - top - bottom_margin
        # Aim for ~0.35" per row, but clamp to fit slide.
        desired_height = 0.35 * rows
        table_height = Inches(min(desired_height, max_table_height_in))

        shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
        table = shape.table

        # Fill cells; convert None to "" to avoid showing "None"
        for r in range(rows):
            for c in range(cols):
                value = self.data_sheet.cell(r + 1, c + 1).value
                table.cell(r, c).text = "" if value is None else str(value)

        # Footnote under table
        footnote_top = table_top + table_height + Inches(0.2)
        footnote_height = Inches(0.4)

        footnote = slide.shapes.add_textbox(table_left, footnote_top, table_width, footnote_height)
        tf = footnote.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "* The values are estimated."
        p.font.size = Pt(max(8, self.settings.text_size_pt - 2))


# Usage example
if __name__ == "__main__":
    converter = ExcelToPPT("data.xlsx", "presentation.pptx")
    converter.create_presentation()